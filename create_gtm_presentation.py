import os
from pptx import Presentation
from pptx.util import Inches

def parse_markdown(md_content):
    """Parses the markdown content to extract sections."""
    sections = []
    current_section = None
    for line in md_content.split('\n'):
        if line.startswith('##'):
            if current_section:
                sections.append(current_section)
            current_section = {'title': line.strip('# ').strip(), 'content': []}
        elif line.startswith('*') and current_section:
            current_section['content'].append(line.strip('* ').strip())
    if current_section:
        sections.append(current_section)
    return sections

def create_presentation(sections, output_filename):
    """Creates a PowerPoint presentation from the parsed sections."""
    prs = Presentation()
    
    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Go-To-Market Strategy"
    subtitle.text = "Senior Care Platform"

    # Content Slides
    for section in sections:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = section['title']
        
        tf = content.text_frame
        tf.clear()
        
        for item in section['content']:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0

    prs.save(output_filename)
    print(f"Presentation '{output_filename}' created successfully.")

if __name__ == '__main__':
    markdown_file = os.path.join('obsidian-vault', 'Strategy', 'Go-To-Market-Strategy.md')
    output_file = 'Go-To-Market-Strategy.pptx'

    try:
        with open(markdown_file, 'r') as f:
            md_content = f.read()
        
        parsed_sections = parse_markdown(md_content)
        create_presentation(parsed_sections, output_file)

    except FileNotFoundError:
        print(f"Error: The file '{markdown_file}' was not found.")
    except Exception as e:
        print(f"An error occurred: {e}")
